from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE

# Utilidades
def add_label(slide, text, left, top, width, height, font_size=16):
    shape = slide.shapes.add_textbox(left, top, width, height)
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(28, 28, 28)
    return shape

def add_input(slide, left, top, width, height, placeholder="", font_size=16):
    # Caja de texto con borde para simular campo editable
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    # Relleno claro
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 250, 252)
    # Borde gris
    line = shape.line
    line.color.rgb = RGBColor(200, 205, 210)
    line.width = Pt(1.25)
    # Texto editable
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = placeholder
    p.font.size = Pt(font_size)
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(100, 100, 100)
    return shape

def add_dropdown(slide, left, top, width, height, placeholder=""):
    shape = add_input(slide, left, top, width, height, placeholder)
    # Triángulo indicando desplegable (flecha hacia abajo)
    tri_w = Inches(0.25)
    tri_h = Inches(0.18)
    tri_left = left + width - tri_w - Inches(0.08)
    tri_top = top + (height - tri_h) / 2
    triangle = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, tri_left, tri_top, tri_w, tri_h)
    triangle.fill.solid()
    triangle.fill.fore_color.rgb = RGBColor(140, 140, 140)
    triangle.line.fill.background()  # sin borde
    return shape

def add_radio_group_si_no(slide, left, top):
    # Círculos para Sí / No
    r = Inches(0.22)
    gap = Inches(0.25)
    # "Sí"
    circ_si = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, r, r)
    circ_si.fill.background()
    circ_si.line.color.rgb = RGBColor(120, 120, 120)
    lbl_si = slide.shapes.add_textbox(left + r + Inches(0.08), top - Inches(0.02), Inches(0.6), Inches(0.3))
    p1 = lbl_si.text_frame.paragraphs[0]
    p1.text = "Sí"
    p1.font.size = Pt(16)
    p1.font.name = "Calibri"
    # "No"
    left_no = left + r + Inches(0.08) + Inches(0.6) + gap
    circ_no = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left_no, top, r, r)
    circ_no.fill.background()
    circ_no.line.color.rgb = RGBColor(120, 120, 120)
    lbl_no = slide.shapes.add_textbox(left_no + r + Inches(0.08), top - Inches(0.02), Inches(0.8), Inches(0.3))
    p2 = lbl_no.text_frame.paragraphs[0]
    p2.text = "No"
    p2.font.size = Pt(16)
    p2.font.name = "Calibri"

def main():
    prs = Presentation()
    # Diapositiva en blanco
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Márgenes y medidas
    left_col = Inches(0.7)
    right_col = Inches(4.1)
    label_w = Inches(3.2)
    field_w = Inches(5.4)
    row_h = Inches(0.55)
    y = Inches(0.6)

    # Título
    title = slide.shapes.add_textbox(Inches(0.7), y, Inches(9), Inches(1))
    tp = title.text_frame.paragraphs[0]
    tp.text = "DATOS CONVOCATORIAS"
    tp.font.bold = True
    tp.font.size = Pt(40)
    tp.font.name = "Calibri"
    tp.alignment = PP_ALIGN.LEFT
    y += Inches(1.0)

    # Función para filas estándar
    def row(label, kind="text", placeholder=""):
        nonlocal y
        add_label(slide, label, left_col, y, label_w, row_h)
        if kind == "text":
            add_input(slide, right_col, y, field_w, row_h, placeholder)
        elif kind == "dropdown":
            add_dropdown(slide, right_col, y, field_w, row_h, placeholder or "Seleccionar…")
        elif kind == "radiosi_no":
            add_radio_group_si_no(slide, right_col, y + Inches(0.1))
        elif kind == "double":
            # Dos campos con guion intermedio
            w = (field_w - Inches(0.5)) / 2
            add_input(slide, right_col, y, w, row_h, placeholder if placeholder else "")
            dash = slide.shapes.add_textbox(right_col + w + Inches(0.2), y, Inches(0.1), row_h)
            dp = dash.text_frame.paragraphs[0]
            dp.text = "-"
            dp.font.size = Pt(18)
            dp.alignment = PP_ALIGN.CENTER
            add_input(slide, right_col + w + Inches(0.4), y, w, row_h, placeholder if placeholder else "")
        y += row_h + Inches(0.18)

    # Campos (texto de la imagen)
    row("Curso Académico")
    row("Código Convocatoria/Programa/Proyecto")
    row("Denominación Convocatoria/Programa/Proyecto")
    row("Plazo de solicitudes: Desde… Hasta", kind="double")
    row("Tipo de entidad contraparte", kind="dropdown", placeholder="Admón europea")
    row("Entidad contraparte (nombre)")
    row("Duración de la estancia: número de meses")
    row("Fechas de estancia: Desde… Hasta", kind="double")
    row("Fuente de Financiación (UA, externa o mixta)")
    row("Subvención:", kind="radiosi_no")
    row("Entidad que subvenciona", kind="dropdown", placeholder="GVA")
    row("Importe convocatoria")

    # Descripción (campo alto)
    add_label(slide, "Descripción de la actividad realizada", left_col, y, label_w, row_h)
    add_input(slide, right_col, y, field_w, Inches(1.2), "")
    y += Inches(1.2) + Inches(0.18)

    row("Publicación", kind="dropdown", placeholder="BOUA")
    row("País de la intervención")
    row("Continente")
    row("Servicio-centro-departamento de gestión en la UA", kind="dropdown", placeholder="Seleccionar…")

    prs.save("datos_convocatorias.pptx")
    print("Generado: datos_convocatorias.pptx")

if __name__ == "__main__":
    main()